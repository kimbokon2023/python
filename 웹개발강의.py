from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# 슬라이드 데이터: 제목과 본문 텍스트
slides_data = [
    ("강의 시작 안내", "강의 제목: '웹 개발의 첫걸음'\n- 강사 소개, 강의 목표 간단 소개\n- 오늘 배울 내용 요약:\n  - 웹의 동작 원리\n  - HTML/CSS/JS의 역할\n  - 개발도구 설치 및 준비\n  - 첫 번째 웹페이지 만들기"),
    ("웹이란 무엇인가?", "- 웹(Web)은 World Wide Web의 줄임말\n- 인터넷 상에서 문서(HTML)와 정보를 주고받는 시스템\n- 브라우저를 통해 접근 가능 (크롬, 엣지 등)"),
    ("웹의 동작 구조", "- 클라이언트(브라우저) → 서버 → 응답 구조\n- HTTP 요청/응답 흐름\n- 정적 페이지 vs 동적 페이지 차이"),
    ("웹 구성 요소", "- HTML: 구조 (뼈대)\n- CSS: 디자인 (옷)\n- JavaScript: 동작과 상호작용 (근육)\n- 이 세 가지가 조화를 이뤄 하나의 웹페이지를 구성"),
    ("웹개발에 필요한 도구들", "- 텍스트 편집기: Notepad++, VSCode 등\n- 웹 브라우저: Chrome (개발자도구 활용 가능)\n- 폴더 관리 및 브라우저 미리보기"),
    ("Notepad++ 설치 방법", "- 공식 홈페이지 방문: https://notepad-plus-plus.org/\n- 설치 파일 다운로드 및 설치 과정 스크린샷\n- 한글 언어 설정 방법 안내"),
    ("Notepad++ 기본 인터페이스", "- 새 파일 생성 방법\n- 탭 구조\n- 줄 번호 보기, 자동 들여쓰기 설정\n- 저장 시 `.html` 확장자로 저장"),
    ("첫 HTML 파일 만들기", "실습: `hello.html` 파일 생성\n```html\n<!DOCTYPE html>\n<html>\n<head>\n  <meta charset='UTF-8'>\n  <title>Hello Web</title>\n</head>\n<body>\n  <h1>Hello, Web!</h1>\n  <p>처음으로 만든 나의 웹페이지입니다.</p>\n</body>\n</html>\n```"),
    ("웹 브라우저에서 열기", "- 파일을 더블 클릭하거나, 브라우저에 드래그하여 열기\n- 코드가 어떻게 표현되는지 확인\n- 간단한 색상, 폰트 수정 실습 안내 (다음 시간 예고)"),
    ("개발자 도구 소개", "- F12 키 또는 우클릭 → 검사\n- HTML 구조 확인\n- Console에서 오류 확인 및 JS 테스트"),
    ("오늘 배운 내용 정리", "- 웹의 구조와 역할 이해\n- Notepad++ 설치 및 사용법 익힘\n- HTML 첫 문서 실습 완료"),
    ("다음 시간 예고", "- HTML 기본 태그들 (제목, 문단, 줄바꿈, 주석 등)\n- 나만의 자기소개 페이지 만들기\n- 미리 준비: 좋아하는 색상, 취미, 소개글 메모해오기"),
    ("질의응답 & 마무리", "- Q&A\n- 과제 안내 (선택): `hello.html` 예제에 글자 색상, 본인 이름 넣어보기\n- 강의자료 제공 예정 안내")
]

# 새 프레젠테이션 생성
prs = Presentation()

# 디자인 적용용 레이아웃 선택 (제목 + 콘텐츠)
layout = prs.slide_layouts[1]

for title, content in slides_data:
    slide = prs.slides.add_slide(layout)
    title_placeholder = slide.shapes.title
    content_placeholder = slide.placeholders[1]

    # 제목 설정
    title_placeholder.text = title

    # 본문 텍스트 설정
    tf = content_placeholder.text_frame
    tf.clear()
    for line in content.split('\n'):
        p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(18)
        p.font.name = '맑은 고딕'
        p.space_after = Pt(4)

# 파일 저장
pptx_path = "1주차_웹개발_Part1.pptx"
prs.save(pptx_path)

pptx_path
